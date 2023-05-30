from selenium import webdriver
from selenium.webdriver.firefox.options import Options
from time import sleep
from pptx import Presentation

# Set up Firefox options to run headless (without opening a browser window)
firefox_options = Options()
firefox_options.headless = True

# Create a new Firefox driver instance
driver = webdriver.Firefox(options=firefox_options)

# Load the webpage
url = 'https://motherson-factsheet.vercel.app/generate-factsheet/21348642'
driver.get(url)
sleep(1)

# Capture a screenshot of the webpage
screenshot_path = "website-screenshot.png"
driver.save_screenshot(screenshot_path)

# Quit the driver
driver.quit()
print("end...")

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a slide with the screenshot
slide_layout = presentation.slide_layouts[6]  # Use the layout for a blank slide
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.add_picture(screenshot_path, left=0, top=0, width=presentation.slide_width, height=presentation.slide_height)

# Save the PowerPoint presentation
presentation.save("factsheet_motherson.pptx")
